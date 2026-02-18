#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR


def create_fashion_presentation():
    # 创建一个新的演示文稿
    prs = Presentation()

    # 设置幻灯片大小为宽屏
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 标题页
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # 标题幻灯片布局
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "20多岁女生过年穿搭指南"
    subtitle.text = "经典耐看款式推荐"

    # 修改标题样式
    title_frame = title.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(139, 0, 0)  # 深红色

    # 副标题样式
    subtitle_para = subtitle.text_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.italic = True

    # 目录页
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容布局
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "目录"
    
    # 添加目录项
    content_text = """1. 红色系单品（新年必备）
2. 经典外套类
3. 内搭上衣
4. 下装推荐
5. 连衣裙系列
6. 配饰点缀
7. 搭配原则
8. 预算分配建议
9. 购买建议"""
    
    content.text = content_text
    
    # 设置字体大小
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)

    # 红色系单品页面
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "红色系单品（新年必备）"
    content_text = """• 正红色毛衣：喜庆又显白，搭配黑色或白色打底裤都很美
• 酒红色连衣裙：优雅大气，适合家庭聚会和走亲访友
• 樱桃红大衣：保暖又有气场，内搭米色高领毛衣更显气质"""

    content.text = content_text
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)

    # 经典外套页面
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "经典外套类"
    content_text = """• 驼色大衣：百搭不挑人，显瘦又高级
• 格纹西装外套：正式场合和日常都能穿，提升精气神
• 短款羽绒服：保暖实用，选择简约款式可穿多年不过时
• 毛呢外套：质感好，版型挺括，显得很有精神"""

    content.text = content_text
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)

    # 内搭上衣页面
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "内搭上衣"
    content_text = """• 高领毛衣：保暖实用，颜色以米色、灰色、黑色为主
• 丝质衬衫：可单穿也可做内搭，适合正式场合
• 针织开衫：温柔知性，可以叠穿增加层次感
• 卫衣：舒适休闲，选择质量好的纯棉款更耐穿"""

    content.text = content_text
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)

    # 下装推荐页面
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "下装推荐"
    content_text = """• 黑色直筒裤：显瘦百搭，职场日常都适用
• 牛仔裤：经典不过时，选择高腰款显腿长
• A字半身裙：优雅减龄，适合各种身材
• 阔腿裤：遮肉显瘦，搭配高跟鞋更显腿长"""

    content.text = content_text
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)

    # 连衣裙系列页面
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "连衣裙系列"
    content_text = """• 针织连衣裙：保暖舒适，适合冬季穿着
• 丝质连衣裙：适合正式场合，显得优雅大方
• 毛呢连衣裙：有质感，不易起皱，适合外出拜访"""

    content.text = content_text
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)

    # 配饰点缀页面
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "配饰点缀"
    content_text = """• 围巾：羊绒材质，颜色以素色为主
• 包包：黑色、棕色等经典色手提包或斜挎包
• 鞋子：小跟靴子、小白鞋、乐福鞋都是不错选择
• 首饰：简约款项链、耳环提升整体精致度"""

    content.text = content_text
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)

    # 搭配原则页面
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "搭配原则"
    content_text = """1. 色彩搭配：主色调不超过三种，保持整体协调
2. 材质选择：注重面料品质，避免廉价感
3. 版型合身：不宜过紧或过松，展现身材曲线
4. 经典款优先：投资一些经典款，可以穿很多年
5. 细节处理：注意衣物整洁，无褶皱、无破损"""

    content.text = content_text
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)

    # 预算分配建议页面
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "预算分配建议"
    content_text = """• 外套：占总预算的40%
• 内搭：占总预算的30%
• 配饰：占总预算的20%
• 特殊场合单品：占总预算的10%"""

    content.text = content_text
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)

    # 购买建议页面
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "购买建议"
    content_text = """• 选择质量好的品牌基础款
• 关注面料成分，优选天然纤维
• 尽量试穿后再购买，确保合身
• 可以选择一些快时尚品牌的质量款作为补充"""

    content.text = content_text
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)

    # 结尾页
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # 标题页布局
    title = slide.shapes.title
    title.text = "谢谢观看！\n\n祝您新年快乐，美丽动人！"
    
    title_frame = title.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(36)
    title_para.font.color.rgb = RGBColor(139, 0, 0)  # 深红色

    # 保存演示文稿
    prs.save('/root/.openclaw/workspace/过年女生穿搭建议.pptx')
    print("PPT文件已成功创建：/root/.openclaw/workspace/过年女生穿搭建议.pptx")


if __name__ == "__main__":
    create_fashion_presentation()